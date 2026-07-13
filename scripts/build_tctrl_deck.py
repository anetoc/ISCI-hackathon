#!/usr/bin/env python3
"""Build the English, judge-facing T-CTRL deck from committed project figures."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "outputs" / "tctrl_hackathon_deck.pptx"

BG = RGBColor(248, 247, 242)
NAVY = RGBColor(16, 52, 70)
TEAL = RGBColor(0, 126, 121)
MINT = RGBColor(218, 241, 235)
GOLD = RGBColor(218, 159, 42)
PALE_GOLD = RGBColor(250, 239, 204)
CORAL = RGBColor(201, 74, 61)
PALE_CORAL = RGBColor(250, 226, 222)
INK = RGBColor(31, 43, 48)
MUTED = RGBColor(91, 106, 111)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(216, 220, 215)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: RGBColor = INK,
    bold: bool = False,
    font: str = "Aptos",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.03,
):
    """Add one text block with explicit typography so rendering is deterministic."""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_header(slide, eyebrow: str, title: str, *, number: int, subtitle: str | None = None):
    """Apply the common release identity and hierarchy to every slide."""

    add_text(slide, eyebrow.upper(), 0.52, 0.30, 6.8, 0.25, size=8.5, color=TEAL, bold=True)
    add_text(slide, title, 0.52, 0.62, 12.1, 0.58, size=25, color=NAVY, bold=True)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(1.25), Inches(12.28), Inches(0.018)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.34, 12.0, 0.42, size=10.5, color=MUTED)
    add_footer(slide, number)


def add_footer(slide, number: int, source: str | None = None):
    """Keep source attribution visible without competing with the scientific message."""

    add_text(
        slide,
        "T-CTRL · ISCI method · Built with Claude: Life Sciences",
        0.52,
        7.17,
        7.1,
        0.16,
        size=6.5,
        color=MUTED,
    )
    if source:
        add_text(slide, source, 7.2, 7.17, 5.35, 0.16, size=6.2, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(
        slide, str(number), 12.58, 7.17, 0.2, 0.16, size=6.5, color=MUTED, align=PP_ALIGN.RIGHT
    )


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LINE,
    radius=True,
):
    """Draw a quiet card that groups one scientific idea."""

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(0.8)
    return card


def add_image_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    """Place a figure without cropping labels or confidence intervals."""

    add_card(slide, x, y, w, h)
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min((w - 0.18) / image_w, (h - 0.18) / image_h)
    draw_w = image_w * scale
    draw_h = image_h * scale
    left = x + (w - draw_w) / 2
    top = y + (h - draw_h) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(draw_w), Inches(draw_h))


def add_metric(
    slide, value: str, label: str, x: float, y: float, w: float, *, color: RGBColor = TEAL
):
    add_text(slide, value, x, y, w, 0.42, size=22, color=color, bold=True)
    add_text(slide, label, x, y + 0.44, w, 0.62, size=9.2, color=INK)


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = BG
    return slide


def build() -> Path:
    """Assemble ten slides from the frozen result hierarchy and committed figures."""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "T-CTRL — Hackathon judge deck (ISCI method)"
    prs.core_properties.subject = "Auditable T-cell controllability analysis"
    prs.core_properties.author = "Abel Costa"
    prs.core_properties.last_modified_by = "Abel Costa"
    prs.core_properties.comments = (
        "Generated by scripts/build_tctrl_deck.py from committed evidence."
    )

    # 1 — Cover
    slide = blank_slide(prs)
    add_text(
        slide,
        "T-CTRL · ISCI METHOD · RESEARCHER TRACK",
        0.58,
        0.46,
        6.0,
        0.25,
        size=8.5,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        "A large cellular change is not biological control.",
        0.58,
        1.05,
        5.85,
        1.52,
        size=31,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        "An auditable system for separating effect size, functional direction and donor-level repeatability in CD4+ T-cell Perturb-seq.",
        0.62,
        2.78,
        5.55,
        1.02,
        size=15,
        color=INK,
    )
    add_card(slide, 0.62, 4.15, 5.42, 1.05, fill=MINT, line=MINT)
    add_text(
        slide,
        "The deliverable is not a target list. It is a method that knows when evidence deserves PASS — and when it must stop.",
        0.88,
        4.42,
        4.92,
        0.55,
        size=12,
        color=NAVY,
        bold=True,
    )
    add_image_contain(slide, ROOT / "figures" / "hackathon_hero.png", 6.65, 0.62, 6.05, 5.95)
    add_text(
        slide,
        "Abel Costa · clinical research and hematology",
        0.62,
        6.57,
        5.4,
        0.25,
        size=8.5,
        color=MUTED,
    )
    add_footer(slide, 1, "Sources: reports/result_lock.md · reports/CLAIM_LEDGER.md")

    # 2 — Clinical question
    slide = blank_slide(prs)
    add_header(
        slide, "The clinical question", "A large effect can still be the wrong answer.", number=2
    )
    concepts = [
        (
            "1",
            "REACH",
            "How much did the cell change?",
            "The size of an abnormal laboratory result.",
        ),
        (
            "2",
            "PRECISION",
            "Which state program moved?",
            "Direction toward activation, Th1/Th2, memory, Treg or exhaustion.",
        ),
        (
            "3",
            "REPEATABILITY",
            "Did independent units agree?",
            "Different guides and donors should reproduce the direction.",
        ),
    ]
    colors = [GOLD, TEAL, NAVY]
    for index, ((number, name, question, analogy), color) in enumerate(
        zip(concepts, colors, strict=True)
    ):
        x = 0.62 + index * 4.18
        add_card(slide, x, 1.65, 3.78, 3.55)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.24), Inches(1.92), Inches(0.48), Inches(0.48)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(
            slide,
            number,
            x + 0.24,
            2.02,
            0.48,
            0.18,
            size=10,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, name, x + 0.86, 1.96, 2.5, 0.32, size=12, color=color, bold=True)
        add_text(slide, question, x + 0.28, 2.72, 3.20, 0.60, size=16, color=NAVY, bold=True)
        add_text(slide, analogy, x + 0.28, 3.58, 3.18, 0.92, size=11, color=INK)
    add_card(slide, 1.22, 5.57, 10.90, 0.78, fill=NAVY, line=NAVY)
    add_text(
        slide,
        "A candidate controller needs all three answers. Magnitude alone is necessary, but not sufficient.",
        1.48,
        5.80,
        10.38,
        0.30,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 2, "Method: docs/method.md")

    # 3 — Experimental reading
    slide = blank_slide(prs)
    add_header(
        slide,
        "How the study reads a perturbation",
        "The useful question is: where did the T cell go?",
        number=3,
        subtitle="Perturb one gene, read the transcriptome, then compare functional direction across guides, donors and contexts.",
    )
    stages = [
        ("CRISPR", "Reduce one defined gene."),
        ("PRIMARY CD4+ T CELL", "Rest and stimulation provide context."),
        ("PERTURB-SEQ", "Thousands of transcripts reveal the state shift."),
        ("AUDITED COMPARISON", "Guide × donor × context tests repetition."),
    ]
    for index, (name, body) in enumerate(stages):
        x = 0.55 + index * 3.18
        fill = MINT if index == 3 else WHITE
        add_card(slide, x, 2.02, 2.72, 1.62, fill=fill)
        add_text(slide, str(index + 1), x + 0.18, 2.20, 0.26, 0.20, size=9, color=TEAL, bold=True)
        add_text(slide, name, x + 0.18, 2.52, 2.35, 0.36, size=11.2, color=NAVY, bold=True)
        add_text(slide, body, x + 0.18, 3.00, 2.32, 0.46, size=9.2, color=INK)
        if index < 3:
            add_text(
                slide,
                "→",
                x + 2.76,
                2.62,
                0.34,
                0.34,
                size=20,
                color=TEAL,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    add_metric(slide, "M", "Magnitude — how far the profile moved", 1.05, 4.50, 2.9, color=GOLD)
    add_metric(
        slide, "C₁", "Precision — movement along a frozen state axis", 5.02, 4.50, 3.2, color=TEAL
    )
    add_metric(
        slide,
        "C₂",
        "Repeatability — agreement across biological units",
        9.30,
        4.50,
        3.0,
        color=NAVY,
    )
    add_card(slide, 1.42, 5.74, 10.48, 0.66, fill=MINT, line=MINT)
    add_text(
        slide,
        "ISCI asks whether C adds information after M is known: the prespecified M → M+C test.",
        1.64,
        5.94,
        10.03,
        0.25,
        size=13,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 3, "Axes: config/axes.yaml · Feature contract: docs/method.md")

    # 4 — Primary result
    slide = blank_slide(prs)
    add_header(
        slide,
        "Primary result",
        "Direction and repeatability add signal beyond magnitude.",
        number=4,
    )
    add_image_contain(slide, ROOT / "figures" / "final" / "fig_central.png", 0.55, 1.55, 7.74, 4.88)
    add_metric(
        slide,
        "+0.357",
        "AUPRC gain in the authoritative full-sample M→M+C test\n95% CI [+0.117,+0.538]",
        8.62,
        1.72,
        3.82,
    )
    add_metric(
        slide,
        "+0.215",
        "Leakage-free grouped OOF gain\n95% CI [+0.074,+0.560] · p=0.010",
        8.62,
        3.15,
        3.82,
        color=NAVY,
    )
    add_metric(
        slide, "0.415→0.722", "Separate descriptive ranking view", 8.62, 4.57, 3.82, color=GOLD
    )
    add_card(slide, 8.58, 5.72, 4.12, 0.70, fill=MINT, line=MINT)
    add_text(
        slide,
        "These are distinct estimands. We never average or silently swap them.",
        8.82,
        5.90,
        3.66,
        0.34,
        size=10.3,
        color=NAVY,
        bold=True,
    )
    add_footer(slide, 4, "Source: reports/result_lock.md · outputs/isci_oof_incremental.json")

    # 5 — Stress tests
    slide = blank_slide(prs)
    add_header(
        slide,
        "Tests designed to break the result",
        "The signal is real — and biologically narrow.",
        number=5,
    )
    add_image_contain(
        slide, ROOT / "figures" / "positive_set_stress_test.png", 0.52, 1.53, 6.08, 3.62
    )
    add_image_contain(
        slide, ROOT / "figures" / "robustness_ablation_gate.png", 6.74, 1.53, 6.08, 3.62
    )
    add_card(slide, 0.72, 5.38, 5.68, 0.86, fill=PALE_CORAL, line=PALE_CORAL)
    add_text(
        slide,
        "External functional-regulator set: FAIL\nΔAUPRC −0.281 [−0.476,−0.073]",
        0.98,
        5.59,
        5.16,
        0.46,
        size=12.2,
        color=CORAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_card(slide, 6.94, 5.38, 5.68, 0.86, fill=PALE_GOLD, line=PALE_GOLD)
    add_text(
        slide,
        "Removing GATA3, TBX21, STAT6 and IRF1 reduces the point gain to near zero (n=9).",
        7.20,
        5.56,
        5.16,
        0.52,
        size=11.2,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 5, "Sources: positive_set_stress_test.json · robustness ablation")

    # 6 — Cross-system scope
    slide = blank_slide(prs)
    add_header(
        slide,
        "Boundary across systems",
        "The property is immune-contextual, not universal.",
        number=6,
    )
    add_image_contain(
        slide,
        ROOT / "figures" / "tctrl_scope_4systems.png",
        0.52,
        1.52,
        8.35,
        4.92,
    )
    scope = [
        ("MARSON CD4+", "PASS · +0.229 [0.072,0.405]", TEAL),
        (
            "SCHMIDT CD4+ / THP-1",
            "Same direction; intervals include zero. Near-miss is not PASS.",
            GOLD,
        ),
        ("K562 / RPE1", "Primary gates fail in non-immune systems.", CORAL),
    ]
    for index, (name, body, color) in enumerate(scope):
        y = 1.72 + index * 1.30
        add_text(slide, name, 9.12, y, 3.28, 0.28, size=10, color=color, bold=True)
        add_text(slide, body, 9.12, y + 0.35, 3.28, 0.66, size=10.5, color=INK)
    add_card(slide, 9.08, 5.63, 3.54, 0.78, fill=NAVY, line=NAVY)
    add_text(
        slide,
        "A boundary is also a result. We do not use the word invariant.",
        9.34,
        5.82,
        3.02,
        0.40,
        size=10.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 6, "Source: figures/tctrl_scope_4systems.png · CLAIM_LEDGER #2")

    # 7 — Translation limits
    slide = blank_slide(prs)
    add_header(
        slide,
        "Translation limits",
        "Controller ≠ therapeutic target ≠ clinical biomarker.",
        number=7,
    )
    add_image_contain(
        slide, ROOT / "figures" / "controller_convergence_quadrant.png", 0.52, 1.50, 6.15, 4.42
    )
    add_image_contain(
        slide, ROOT / "outputs" / "iec_clinical" / "iec_prediction.png", 6.83, 1.50, 5.98, 4.42
    )
    add_card(slide, 0.72, 6.06, 5.74, 0.62, fill=PALE_GOLD, line=PALE_GOLD)
    add_text(
        slide,
        "IRF1 ranks #1 as a controller but points in an unfavorable intervention direction.",
        0.96,
        6.22,
        5.26,
        0.30,
        size=10.4,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_card(slide, 7.03, 6.06, 5.58, 0.62, fill=PALE_CORAL, line=PALE_CORAL)
    add_text(
        slide,
        "CAR-T response: NULL · study-out AUROC 0.533 · CD8 baseline 0.585 · 87 patients / 9 studies",
        7.26,
        6.18,
        5.12,
        0.38,
        size=9.9,
        color=CORAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 7, "Sources: controller_convergence_quadrant.png · iec_clinical/verdict.json")

    # 8 — Judgment product
    slide = blank_slide(prs)
    add_header(
        slide, "The innovation is judgment", "The product also knows when to stop.", number=8
    )
    add_image_contain(slide, ROOT / "figures" / "layer_verdict_map.png", 0.52, 1.54, 8.52, 4.90)
    verdicts = [
        ("PASS", "Prespecified direction and every required gate pass.", TEAL, MINT),
        ("FAIL", "The claim is testable and a required gate fails.", CORAL, PALE_CORAL),
        ("NULL", "The test is valid but adds no supported signal.", NAVY, RGBColor(227, 235, 239)),
        ("NOT-EVALUABLE", "A required input is absent; no biology is forced.", GOLD, PALE_GOLD),
    ]
    for index, (name, body, color, fill) in enumerate(verdicts):
        y = 1.58 + index * 1.18
        add_card(slide, 9.30, y, 3.34, 0.96, fill=fill, line=fill)
        add_text(slide, name, 9.52, y + 0.14, 2.90, 0.22, size=10.5, color=color, bold=True)
        add_text(slide, body, 9.52, y + 0.44, 2.90, 0.38, size=8.8, color=INK)
    add_card(slide, 1.22, 6.55, 10.90, 0.38, fill=NAVY, line=NAVY)
    add_text(
        slide,
        "Strong science does not maximize positive results; it makes every conclusion auditable and falsifiable.",
        1.48,
        6.64,
        10.38,
        0.18,
        size=10.3,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 8, "Source: outputs/hackathon/claim_manifest.json")

    # 9 — Prospective falsification
    slide = blank_slide(prs)
    add_header(
        slide,
        "Gladstone bridge",
        "A prospective experiment that can falsify the next claim.",
        number=9,
    )
    metrics = [
        ("25", "target genes"),
        ("54", "CRISPR guides"),
        ("8–12", "independent donors"),
        ("2", "paired contexts"),
    ]
    for index, (value, label) in enumerate(metrics):
        x = 0.62 + index * 3.04
        add_card(slide, x, 1.62, 2.62, 1.20, fill=MINT if index < 2 else PALE_GOLD)
        add_text(
            slide,
            value,
            x + 0.18,
            1.84,
            2.26,
            0.38,
            size=22,
            color=TEAL if index < 2 else GOLD,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide, label, x + 0.18, 2.28, 2.26, 0.24, size=9.4, color=INK, align=PP_ALIGN.CENTER
        )
    add_card(slide, 0.62, 3.18, 7.64, 2.48)
    add_text(slide, "DONOR-RESOLVED DESIGN", 0.92, 3.46, 3.4, 0.28, size=10, color=TEAL, bold=True)
    add_text(
        slide,
        "The donor is the biological replicate. Guides, wells and sequencing channels do not increase n. Each donor contributes matched rest and stimulation contexts, and promotion requires concordant direction across the predeclared panel.",
        0.92,
        3.92,
        6.98,
        1.08,
        size=14,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        "Falsification: a residual precision gain below the frozen lower bound, donor-bootstrap interval crossing zero, or gene-label/context exchange p≥0.05 blocks promotion.",
        0.92,
        5.06,
        6.98,
        0.40,
        size=9.4,
        color=MUTED,
    )
    add_card(slide, 8.58, 3.18, 4.18, 2.48, fill=PALE_CORAL, line=PALE_CORAL)
    add_text(
        slide, "STOP BEFORE SYNTHESIS", 8.88, 3.48, 3.58, 0.30, size=11, color=CORAL, bold=True
    )
    add_text(
        slide,
        "Reference sequence, vector compatibility, on-target activity and off-target risk must pass first. PAPOLG-1 remains low-support. A proposed panel is not an experimentally validated result.",
        8.88,
        4.02,
        3.58,
        1.12,
        size=12.2,
        color=INK,
        bold=True,
    )
    add_footer(slide, 9, "Source: reports/PROSPECTIVE_DONOR_PANEL_PROTOCOL.md")

    # 10 — Delivery
    slide = blank_slide(prs)
    add_header(
        slide,
        "Open-source delivery",
        "The artifact is executable, reusable and auditable.",
        number=10,
    )
    layers = [
        ("1 · CONTRACTS", "Frozen axes, claims, gates and DatasetSpec", PALE_GOLD),
        ("2 · PYTHON ENGINE", "CLI, adapters, block-bounded H5AD processing", MINT),
        ("3 · EVIDENCE", "JSON/CSV, figures, hashes, tests and claim ledger", WHITE),
        ("4 · EXPERIENCE", "Notebook, offline HTML, Full-HD video and deck", PALE_CORAL),
    ]
    for index, (name, body, fill) in enumerate(layers):
        x = 0.52 + index * 3.15
        add_card(slide, x, 1.66, 2.78, 1.72, fill=fill)
        add_text(
            slide,
            name,
            x + 0.20,
            1.92,
            2.38,
            0.28,
            size=9.5,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, body, x + 0.20, 2.43, 2.38, 0.58, size=10, color=INK, align=PP_ALIGN.CENTER)
    add_card(slide, 0.72, 3.84, 7.64, 1.52, fill=NAVY, line=NAVY)
    add_text(
        slide, "RUN IT ON A NEW DATASET", 1.02, 4.08, 3.2, 0.26, size=9.5, color=MINT, bold=True
    )
    add_text(
        slide,
        "uv run isci pipeline dataset.yaml",
        1.02,
        4.53,
        6.98,
        0.42,
        size=20,
        color=WHITE,
        bold=True,
        font="Aptos Mono",
    )
    add_text(
        slide,
        "Validate → inspect → construct effects → extract features → audit output",
        1.02,
        5.02,
        6.98,
        0.24,
        size=9.5,
        color=WHITE,
    )
    add_card(slide, 8.66, 3.84, 4.02, 1.52, fill=MINT, line=MINT)
    add_metric(
        slide,
        "21/21",
        "automated release gates passed\nHuman rehearsal and PI wording gates remain explicit",
        9.00,
        4.02,
        3.36,
    )
    add_card(slide, 1.22, 5.83, 10.90, 0.70, fill=TEAL, line=TEAL)
    add_text(
        slide,
        "Delivery = scientific software + frozen evidence + visual explanation + an honest stopping rule.",
        1.48,
        6.04,
        10.38,
        0.30,
        size=13.2,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 10, "MIT code · third-party data retain source terms · CITATION.cff")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output = build()
    print(f"Wrote {output}")
