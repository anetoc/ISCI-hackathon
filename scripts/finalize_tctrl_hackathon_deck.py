#!/usr/bin/env python3
"""Finalize the redesigned judge deck without rebuilding its hand-tuned layout.

The public PPTX is intentionally treated as the visual template. This script applies the small,
auditable corrections that must stay synchronized with the public naming and result contract.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_DECK = ROOT / "outputs" / "tctrl_hackathon_deck.pptx"
DEFAULT_FIGURE = ROOT / "figures" / "tctrl_scope_4systems.png"


def replace_public_copy(prs: Presentation) -> int:
    """Remove stale branding and the incorrect shrinkage framing from visible text."""

    replacements = {
        "powered by ISCI": "ISCI method",
        (
            "These are distinct estimands — never averaged or swapped. With only 13 positives "
            "we re-ran everything out-of-fold and report the shrinkage."
        ): "These are two different measurements — not one number shrinking. We report both.",
    }
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    updated = run.text
                    for old, new in replacements.items():
                        updated = updated.replace(old, new)
                    if updated != run.text:
                        run.text = updated
                        changed += 1
    return changed


def replace_scope_figure(prs: Presentation, figure: Path) -> None:
    """Replace slide 6's sole evidence image while preserving its exact geometry."""

    if len(prs.slides) != 10:
        raise ValueError(f"Expected the 10-slide judge deck, found {len(prs.slides)} slides")
    slide = prs.slides[5]
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if len(pictures) != 1:
        raise ValueError(f"Expected one picture on slide 6, found {len(pictures)}")
    old = pictures[0]
    left, top, width, height = old.left, old.top, old.width, old.height
    old._element.getparent().remove(old._element)
    new = slide.shapes.add_picture(str(figure), left, top, width, height)
    new.name = "T-CTRL four-system controllability boundary"
    new._element.nvPicPr.cNvPr.set(
        "descr", "Four-system scope map showing PASS, near-miss and FAIL verdicts."
    )


def finalize(deck: Path, figure: Path) -> Path:
    """Apply the final corrections atomically so a failed save cannot damage the deck."""

    if not deck.is_file():
        raise FileNotFoundError(deck)
    if not figure.is_file():
        raise FileNotFoundError(figure)
    prs = Presentation(deck)
    replace_public_copy(prs)
    replace_scope_figure(prs, figure)
    prs.core_properties.title = "T-CTRL — Which genes actually steer T-cell state?"
    prs.core_properties.subject = "Auditable T-cell controllability analysis using the ISCI method"
    temporary = deck.with_name(f".{deck.stem}.tmp{deck.suffix}")
    prs.save(temporary)
    temporary.replace(deck)
    return deck


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--deck", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--figure", type=Path, default=DEFAULT_FIGURE)
    args = parser.parse_args()
    output = finalize(args.deck.resolve(), args.figure.resolve())
    print(f"Finalized {output}")


if __name__ == "__main__":
    main()
